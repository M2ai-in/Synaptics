import json
import re
import os
from pptx import Presentation
from pptx.util import Pt
from typing import List, Dict, Union
from .base import Tool
class SlideGenerationTool(Tool):
    name: str = "slide_generation_tool"
    description: str = ("A tool that can create a PPTX deck for content. It takes a list of dictionaries. Each dictionary represents a slide with two keys: 'slide_title' and 'content'.")
    arg: str = "List[Dict[slide_title, content]]. Ensure the Action Input is JSON parseable so I can convert it to required format"
    def __init__(self, client_details: dict = None, **data):
        super().__init__(client_details=client_details, **data)
    def run(self, slide_content: Union[str, List[Dict[str, str]]], temp = 0.7, max_tokens= 4000) -> dict:
        print(f"ğŸ“¥ Slide Generation Tool received input of type: {type(slide_content)} with temp: {temp}, max_tokens: {max_tokens}")
        print("ğŸ› ï¸ Processing slide content...")
        if isinstance(slide_content, str):
            try:
                slide_content = json.loads(slide_content)
                print("âœ… Parsed JSON input successfully.")
            except json.JSONDecodeError as e:
                return {"error": f"âŒ Failed to parse input as JSON: {str(e)}", "received_type": str(type(slide_content)), "raw_input": slide_content}
        if not isinstance(slide_content, list) or not all(isinstance(slide, dict) for slide in slide_content):
            return {"error": "âŒ Input must be a list of dictionaries with 'slide_title' and 'content'.", "received_type": str(type(slide_content)), "raw_input": slide_content}
        for i, slide in enumerate(slide_content):
            if "slide_title" not in slide or "content" not in slide:
                return {"error": f"âŒ Slide {i} is missing 'slide_title' or 'content'.", "slide_data": slide}
        presentation = Presentation()
        for i, slide in enumerate(slide_content):
            slide_layout = presentation.slide_layouts[1]
            ppt_slide = presentation.slides.add_slide(slide_layout)
            title_shape = ppt_slide.shapes.title
            title_shape.text = slide["slide_title"]
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            body_shape = ppt_slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            content_lines = slide["content"].split("\n")
            for j, line in enumerate(content_lines):
                if j == 0: p = tf.paragraphs[0]
                else: p = tf.add_paragraph()
                p.text = line.strip()
                p.level = 0
                p.font.size = Pt(20)
        topic_title = slide_content[0]["slide_title"] if slide_content else "Presentation"
        safe_title = re.sub(r'[^0-9a-zA-Z]+', '_', topic_title)
        output_path = f"{safe_title}.pptx"
        presentation.save(output_path)
        print("Saved to ", output_path)
        return {"message": "âœ… Slide deck created successfully!", "file_path": output_path, "slide_count": len(slide_content)}
